import os
import json
import logging
import re
from fastapi import APIRouter, UploadFile, File, BackgroundTasks, HTTPException
from fastapi.responses import JSONResponse, StreamingResponse
from pydantic import BaseModel
from typing import Optional
import shutil

logger = logging.getLogger(__name__)

# 借用已经在根目录下的 ingest 模块
import sys
# 保证能引用到 ingest_master_md
sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
from ingest import ingest_master_md

# 引入现成的 RAGEngine 用于进行在线作答
from app.services.rag_engine import rag_engine

from fastapi import Depends
from app.api.auth import verify_jwt
router = APIRouter()

DATA_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), "data")
os.makedirs(DATA_DIR, exist_ok=True)

class IngestRequest(BaseModel):
    filename: str
    lecture_id: str

class ChatRequest(BaseModel):
    lecture_id: str
    question: str
    image_base64: Optional[str] = None  # 支持可选的图片 Base64 数据

@router.get("/files")
async def list_files():
    """获取 data 目录下的 markdown 文件列表"""
    files = []
    if os.path.exists(DATA_DIR):
        for f in os.listdir(DATA_DIR):
            if f.endswith(".md"):
                # 获取简单的文件信息
                file_path = os.path.join(DATA_DIR, f)
                if f.endswith(".master.md"):
                    namespace = f[:-10]
                else:
                    namespace = f[:-3]
                files.append({
                    "name": f,
                    "size": os.path.getsize(file_path),
                    # 推导命名空间前缀
                    "suggested_namespace": namespace,
                    # 动态探测向量库：该讲义是否已经生成完毕并在内存中提供服务
                    "is_ingested": rag_engine.check_ingested(namespace)
                })
    return {"status": "success", "files": files}

def convert_pdf_to_md(file_path: str) -> str:
    """将PDF转换为Markdown"""
    try:
        import pdfplumber
    except ImportError:
        raise HTTPException(status_code=500, detail="请安装 pdfplumber: pip install pdfplumber")
    
    md_content = ""
    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text:
                md_content += f"## 第{page_num}页\n\n{text}\n\n"
    return md_content

def convert_docx_to_md(file_path: str) -> str:
    """将Word文档转换为Markdown"""
    try:
        from docx import Document
    except ImportError:
        raise HTTPException(status_code=500, detail="请安装 python-docx: pip install python-docx")
    
    md_content = ""
    doc = Document(file_path)
    
    for para in doc.paragraphs:
        if para.style.name.startswith('Heading'):
            level = para.style.name[-1] if para.style.name[-1].isdigit() else '1'
            md_content += f"{'#' * int(level)} {para.text}\n\n"
        else:
            md_content += f"{para.text}\n\n"
    
    for table in doc.tables:
        md_content += "\n| "
        for cell in table.rows[0].cells:
            md_content += cell.text + " | "
        md_content += "\n| "
        for _ in table.rows[0].cells:
            md_content += "--- | "
        md_content += "\n"
        
        for row in table.rows[1:]:
            md_content += "| "
            for cell in row.cells:
                md_content += cell.text + " | "
            md_content += "\n"
        md_content += "\n"
    
    return md_content

def convert_pptx_to_md(file_path: str) -> str:
    """将PPT转换为Markdown"""
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(status_code=500, detail="请安装 python-pptx: pip install python-pptx")
    
    md_content = ""
    prs = Presentation(file_path)
    
    for slide_num, slide in enumerate(prs.slides, 1):
        md_content += f"## 幻灯片 {slide_num}\n\n"
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                md_content += f"{shape.text}\n\n"
        
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            md_content += f"**备注:** {slide.notes_slide.notes_text_frame.text}\n\n"
        
        md_content += "---\n\n"
    
    return md_content

@router.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """接收浏览器上传的文件，支持 .md/.pdf/.docx/.pptx"""
    allowed_extensions = {'.md', '.pdf', '.docx', '.pptx'}
    ext = os.path.splitext(file.filename)[1].lower()
    
    if ext not in allowed_extensions:
        raise HTTPException(status_code=400, detail=f"Only {allowed_extensions} files are allowed")
    
    original_filename = file.filename
    temp_path = os.path.join(DATA_DIR, f"temp_{original_filename}")
    
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    md_filename = os.path.splitext(original_filename)[0] + ".md"
    md_path = os.path.join(DATA_DIR, md_filename)
    
    if ext == '.md':
        shutil.move(temp_path, md_path)
        message = "MD文件上传成功"
    else:
        try:
            if ext == '.pdf':
                md_content = convert_pdf_to_md(temp_path)
            elif ext == '.docx':
                md_content = convert_docx_to_md(temp_path)
            elif ext == '.pptx':
                md_content = convert_pptx_to_md(temp_path)
            else:
                raise HTTPException(status_code=400, detail="Unsupported file type")
            
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(md_content)
            
            os.remove(temp_path)
            message = f"{ext.upper()}文件已转换为Markdown"
        except Exception as e:
            if os.path.exists(temp_path):
                os.remove(temp_path)
            raise HTTPException(status_code=500, detail=f"文件转换失败: {str(e)}")
    
    return {"status": "success", "filename": md_filename, "message": message}

@router.delete("/files/{filename}")
async def delete_file(filename: str):
    """删除 .md 文件"""
    file_path = os.path.join(DATA_DIR, filename)
    if os.path.exists(file_path):
        os.remove(file_path)
        
        # 物理删除了该文档的同时，务必将其早先写入向量库（ChromaDB）的魂魄骨骸一并销毁！
        lecture_id = "lecture_" + filename.replace(".master.md", "")
        try:
            existing = rag_engine.vector_store.get(where={"lecture_id": lecture_id})
            if existing and existing.get("ids"):
                rag_engine.vector_store.delete(ids=existing["ids"])
                print(f"✅ [清道夫] 成功销毁 {filename} 残留在数据库中的灵魂碎片 ({len(existing['ids'])} 条)。")
            else:
                print(f"⚠️ [清道夫] {filename} 的物理文件已删，但向量库暂无记录。")
        except Exception as e:
            print(f"⚠️ [清道夫] 回收向量切片时出错，但不阻塞前台: {e}")
            
        return {"status": "success", "message": "File and corresponding vector chunks deleted"}
    raise HTTPException(status_code=404, detail="File not found")

@router.get("/files/{namespace}/chunks")
async def get_file_chunks(namespace: str):
    """透视特定 namespace 下的所有向量库切片"""
    try:
        chunks = rag_engine.get_namespace_chunks(namespace)
        return {"status": "success", "namespace": namespace, "total": len(chunks), "chunks": chunks}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/ingest")
async def trigger_ingest(request: IngestRequest):
    """触发向量化入库（同步阻塞，便于前端获知确切完成时间）"""
    file_path = os.path.join(DATA_DIR, request.filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found in data directory")
    
    try:
        # FastAPI 会自动在线程池中运行同步函数，不阻塞主事件循环
        ingest_master_md(file_path, request.lecture_id)
        
        return {
            "status": "success", 
            "message": f"{request.filename} 的切片与高维入库已经彻底完成！"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/chat")
async def process_chat(request: ChatRequest):
    """用于测试向量检索引擎的对话接口"""
    try:
        if (not request.question or not request.question.strip()) and not request.image_base64:
            raise HTTPException(status_code=400, detail="Question cannot be empty unless an image is provided")
            
        answer = await rag_engine.get_answer(
            question=request.question, 
            lecture_id=request.lecture_id,
            image_base64=request.image_base64
        )
        return {"status": "success", "answer": answer}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/chat/stream")
async def process_chat_stream(request: ChatRequest):
    """流式 RAG 对话接口"""
    try:
        if (not request.question or not request.question.strip()) and not request.image_base64:
            raise HTTPException(status_code=400, detail="Question cannot be empty unless an image is provided")
            
        return StreamingResponse(
            rag_engine.get_answer_stream(
                question=request.question,
                lecture_id=request.lecture_id,
                image_base64=request.image_base64
            ),
            media_type="text/event-stream"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

class ExamRequest(BaseModel):
    filename: str
    lecture_id: str
    difficulty: str = "medium"
    question_count: int = 5

class GradeRequest(BaseModel):
    questions: list
    answers: list

@router.post("/exam/analyze")
async def analyze_key_points(request: ExamRequest):
    """分析文件并提取重难点"""
    file_path = os.path.join(DATA_DIR, request.filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found in data directory")

    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

        if not content or len(content.strip()) == 0:
            raise HTTPException(status_code=400, detail="File content is empty")

        from app.core.config import settings
        from langchain_openai import ChatOpenAI

        _is_siliconflow = "siliconflow" in settings.LLM_API_BASE.lower()
        _no_think_params = (
            {"enable_thinking": False}
            if _is_siliconflow
            else {"chat_template_kwargs": {"enable_thinking": False}}
        )

        llm = ChatOpenAI(
            openai_api_base=settings.LLM_API_BASE,
            openai_api_key=settings.LLM_API_KEY,
            model_name=settings.LLM_MODEL_NAME,
            temperature=0.3,
            max_tokens=4096,
            extra_body={
                **_no_think_params,
                "stop": ["<|im_end|>", "<|endoftext|>"]
            }
        )

        analysis_prompt = f"""请分析以下教材内容，提取重难点并总结。

要求：
1. 识别章节标题和小节标题
2. 标注每个章节的核心概念和重要术语
3. 指出需要重点理解的知识点
4. 标记可能的考试重点

教材内容：
{content[:15000]}

请以JSON格式输出，包含以下字段：
- "chapters": [{{"title": "章节标题", "key_points": ["要点1", "要点2"], "difficulty": "high/medium/low"}}]
- "important_concepts": [{{"term": "术语", "definition": "定义", "importance": "high/medium/low"}}]
- "summary": "整体总结"

请确保输出是有效的JSON格式。"""

        logger.info(f"Analyzing file: {request.filename}, content length: {len(content)}")
        response = llm.invoke(analysis_prompt)
        logger.info(f"LLM response length: {len(response.content)}")

        try:
            json_match = re.search(r'\{[\s\S]*\}', response.content)
            if json_match:
                analysis_result = json.loads(json_match.group())
            else:
                analysis_result = {
                    "chapters": [],
                    "important_concepts": [],
                    "summary": response.content[:500],
                    "error": "JSON解析失败，返回原始内容"
                }
        except json.JSONDecodeError as je:
            logger.error(f"JSON decode error: {je}")
            analysis_result = {
                "chapters": [],
                "important_concepts": [],
                "summary": response.content[:500],
                "error": f"JSON解析失败: {str(je)}"
            }

        return {"status": "success", "analysis": analysis_result}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"analyze_key_points error: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"分析失败: {str(e)}")

@router.post("/exam/generate")
async def generate_exam(request: ExamRequest):
    """生成考题"""
    file_path = os.path.join(DATA_DIR, request.filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found in data directory")

    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

        if not content or len(content.strip()) == 0:
            raise HTTPException(status_code=400, detail="File content is empty")

        from app.core.config import settings
        from langchain_openai import ChatOpenAI

        _is_siliconflow = "siliconflow" in settings.LLM_API_BASE.lower()
        _no_think_params = (
            {"enable_thinking": False}
            if _is_siliconflow
            else {"chat_template_kwargs": {"enable_thinking": False}}
        )

        llm = ChatOpenAI(
            openai_api_base=settings.LLM_API_BASE,
            openai_api_key=settings.LLM_API_KEY,
            model_name=settings.LLM_MODEL_NAME,
            temperature=0.5,
            max_tokens=4096,
            extra_body={
                **_no_think_params,
                "stop": ["<|im_end|>", "<|endoftext|>"]
            }
        )

        difficulty_desc = {
            "easy": "基础概念题，只需回忆和理解基本定义",
            "medium": "应用题，需要理解知识点并能简单应用",
            "hard": "综合分析题，需要深入理解并能分析和解决问题"
        }

        generation_prompt = f"""请根据以下教材内容生成{request.question_count}道{difficulty_desc.get(request.difficulty, '中等')}难度的考题。

要求：
1. 题目必须基于教材内容，不能超出教材范围
2. 题型包括：选择题（单选/多选）、判断题、填空题、简答题、论述题
3. 每道题都要有明确的正确答案
4. 难度要适中，符合{difficulty_desc.get(request.difficulty, '中等')}水平

教材内容：
{content[:15000]}

请以JSON格式输出，包含以下字段：
- "questions": [
    {{
        "id": 1,
        "type": "choice(单选)/multiple(多选)/true_false(判断)/fill_blank(填空)/short(简答)/essay(论述)",
        "text": "题目内容",
        "question": "题目内容(同text)",
        "options": ["A. 选项1", "B. 选项2", "C. 选项3", "D. 选项4"],  // 选择题必须有
        "answer": "正确答案",
        "reference_answer": "参考答案(详细)",
        "rubric": "评分标准(如：内容完整性30分、逻辑清晰30分、表达准确40分)",  // 主观题必须有
        "analysis": "解题思路",
        "points": 分数,
        "knowledge_tag": "知识点标签(如：第一章-档案学基础)"
    }}
  ]

请确保输出是有效的JSON格式。"""

        logger.info(f"Generating exam: {request.filename}, count: {request.question_count}, difficulty: {request.difficulty}")
        response = llm.invoke(generation_prompt)
        logger.info(f"LLM response length: {len(response.content)}")

        try:
            json_match = re.search(r'\[[\s\S]*\]', response.content)
            if json_match:
                questions = json.loads(json_match.group())
            else:
                questions = [{"error": "JSON解析失败", "raw": response.content[:500]}]
        except json.JSONDecodeError as je:
            logger.error(f"JSON decode error: {je}")
            questions = [{"error": f"JSON解析失败: {str(je)}", "raw": response.content[:500]}]

        return {"status": "success", "questions": questions}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"generate_exam error: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"生成考题失败: {str(e)}")

def parse_json_response(content: str):
    """解析JSON响应，支持markdown代码块"""
    try:
        content = content.strip()
        if content.startswith("```"):
            content = re.sub(r'^```\w*\n?', '', content)
            content = re.sub(r'\n?```$', '', content)
        json_match = re.search(r'\{[\s\S]*\}', content)
        if json_match:
            return json.loads(json_match.group())
    except Exception as e:
        logger.error(f"JSON parse error: {e}")
    return None

def grade_single_choice(q: dict, student_ans: str) -> dict:
    """单选题批改：完全匹配满分，否则0分"""
    correct = str(q.get('answer', '')).strip().upper()
    student = student_ans.strip().upper()
    is_correct = correct == student
    return {
        "question_id": q.get('id', q.get('question_id', 0)),
        "score": q.get('points', q.get('score', 5)) if is_correct else 0,
        "max_score": q.get('points', q.get('score', 5)),
        "is_correct": is_correct,
        "correct_answer": correct,
        "student_answer": student_ans,
        "analysis": "正确" if is_correct else f"正确答案是 {correct}",
        "knowledge_tag": q.get('knowledge_tag', ''),
        "weakness": [] if is_correct else [q.get('knowledge_tag', '')]
    }

def grade_multiple_choice(q: dict, student_ans: str) -> dict:
    """多选题批改：完全匹配满分，漏选按比例给分，多选/错选0分"""
    correct_set = set(str(q.get('answer', '')).upper().split(','))
    student_set = set(student_ans.strip().upper().split(','))
    correct_count = len(correct_set)
    
    if student_set == correct_set:
        score = q.get('points', q.get('score', 5))
        is_correct = True
        analysis = "完全正确"
    elif student_set < correct_set:
        ratio = len(student_set & correct_set) / correct_count
        score = round(ratio * q.get('points', q.get('score', 5)), 1)
        is_correct = False
        analysis = f"部分正确，漏选 {correct_set - student_set}"
    else:
        score = 0
        is_correct = False
        analysis = "多选或错选"
    
    return {
        "question_id": q.get('id', q.get('question_id', 0)),
        "score": score,
        "max_score": q.get('points', q.get('score', 5)),
        "is_correct": is_correct,
        "correct_answer": ','.join(sorted(correct_set)),
        "student_answer": student_ans,
        "analysis": analysis,
        "knowledge_tag": q.get('knowledge_tag', ''),
        "weakness": [] if is_correct else [q.get('knowledge_tag', '')]
    }

def grade_true_false(q: dict, student_ans: str) -> dict:
    """判断题批改：同单选题"""
    return grade_single_choice(q, student_ans)

def grade_fill_blank(q: dict, student_ans: str, fuzzy: bool = True) -> dict:
    """填空题批改：支持模糊匹配"""
    correct = str(q.get('answer', '')).strip()
    student = student_ans.strip()
    
    is_correct = False
    if fuzzy:
        import difflib
        ratio = difflib.SequenceMatcher(None, student, correct).ratio()
        is_correct = ratio >= 0.8 or correct in student or student in correct
    else:
        is_correct = correct == student
    
    return {
        "question_id": q.get('id', q.get('question_id', 0)),
        "score": q.get('points', q.get('score', 5)) if is_correct else 0,
        "max_score": q.get('points', q.get('score', 5)),
        "is_correct": is_correct,
        "correct_answer": correct,
        "student_answer": student_ans,
        "analysis": "正确" if is_correct else f"正确答案：{correct}",
        "knowledge_tag": q.get('knowledge_tag', ''),
        "weakness": [] if is_correct else [q.get('knowledge_tag', '')]
    }

def build_subjective_prompt(q: dict, student_ans: str) -> str:
    """构建主观题评分prompt"""
    rubric = q.get('rubric', '')
    reference = q.get('reference_answer', q.get('reference', ''))
    
    prompt = f"""请批改以下主观题答案。

题目：{q.get('question', q.get('text', ''))}

学生答案：{student_ans}

参考答案：{reference}

评分标准：{rubric}

请以JSON格式输出评分结果，包含以下字段：
{{
    "score": 实际得分,
    "max_score": 满分,
    "dimension_scores": [{{"dimension": "维度名", "score": 得分, "comment": "评语"}}],
    "total_comment": "总体评语",
    "strengths": ["优点1", "优点2"],
    "weaknesses": ["不足1", "不足2"],
    "suggestions": ["建议1", "建议2"]
}}

请确保输出是有效的JSON格式。"""
    return prompt

@router.post("/exam/grade")
async def grade_exam(request: GradeRequest):
    """批改考题 - 支持多种题型"""
    try:
        from app.core.config import settings
        from langchain_openai import ChatOpenAI
        import difflib

        _is_siliconflow = "siliconflow" in settings.LLM_API_BASE.lower()
        _no_think_params = (
            {"enable_thinking": False}
            if _is_siliconflow
            else {"chat_template_kwargs": {"enable_thinking": False}}
        )

        llm = ChatOpenAI(
            openai_api_base=settings.LLM_API_BASE,
            openai_api_key=settings.LLM_API_KEY,
            model_name=settings.LLM_MODEL_NAME,
            temperature=0.1,
            max_tokens=4096,
            extra_body={
                **_no_think_params,
                "stop": ["<|im_end|>", "<|endoftext|>"]
            }
        )

        # 构建题目和答案的映射
        question_map = {str(q.get('id', q.get('question_id', i))): q for i, q in enumerate(request.questions)}
        answer_map = {}
        for a in request.answers:
            qid = str(a.get('qid', a.get('question_id', '')))
            answer_map[qid] = a.get('answer', '')

        # 分批处理：客观题本地批改，主观题LLM批改
        objective_results = []
        subjective_questions = []
        subjective_prompts = []
        
        for i, q in enumerate(request.questions):
            qid = str(q.get('id', q.get('question_id', f'q_{i}')))
            q_type = (q.get('type') or q.get('question_type') or '').lower()
            student_ans = answer_map.get(qid, '')
            
            if q_type in ['choice', 'single', 'single_choice']:
                # 判断是单选还是多选
                if ',' in str(q.get('answer', '')) or '多选' in str(q.get('text', '')):
                    objective_results.append(grade_multiple_choice(q, student_ans))
                else:
                    objective_results.append(grade_single_choice(q, student_ans))
            elif q_type in ['multiple', 'multiple_choice']:
                objective_results.append(grade_multiple_choice(q, student_ans))
            elif q_type in ['true_false', 'judge', 'judgment', '判断']:
                objective_results.append(grade_true_false(q, student_ans))
            elif q_type in ['fill', 'fill_blank', 'blank', '填空']:
                objective_results.append(grade_fill_blank(q, student_ans, fuzzy=True))
            elif q_type in ['short', 'essay', '简答', '论述', '主观']:
                subjective_questions.append((q, student_ans))
                subjective_prompts.append(build_subjective_prompt(q, student_ans))

        # 批量调用LLM批改主观题
        subjective_results = []
        if subjective_prompts:
            logger.info(f"Grading {len(subjective_prompts)} subjective questions with LLM")
            
            for idx, (q, student_ans) in enumerate(subjective_questions):
                try:
                    prompt = build_subjective_prompt(q, student_ans)
                    response = llm.invoke(prompt)
                    parsed = parse_json_response(response.content)
                    
                    if parsed:
                        parsed['question_id'] = q.get('id', f'subj_{idx}')
                        parsed['max_score'] = q.get('points', q.get('score', 10))
                        subjective_results.append(parsed)
                    else:
                        subjective_results.append({
                            "question_id": q.get('id', f'subj_{idx}'),
                            "score": 0,
                            "max_score": q.get('points', q.get('score', 10)),
                            "total_comment": "LLM解析失败",
                            "suggestions": ["请重新提交批改"]
                        })
                except Exception as e:
                    logger.error(f"LLM grading error for q {idx}: {e}")
                    subjective_results.append({
                        "question_id": q.get('id', f'subj_{idx}'),
                        "score": 0,
                        "max_score": q.get('points', q.get('score', 10)),
                        "total_comment": f"批改出错: {str(e)}",
                        "suggestions": ["请重新提交批改"]
                    })

        # 为主观题结果添加question_id
        for i, sr in enumerate(subjective_results):
            if 'question_id' not in sr:
                subj_idx = i
                if subj_idx < len(subjective_questions):
                    sr['question_id'] = subjective_questions[subj_idx][0].get('id', f'subj_{i}')
        
        # 合并结果
        all_results = objective_results + subjective_results
        
        # 计算总分
        total_score = sum(r.get('score', 0) for r in all_results)
        max_total = sum(r.get('max_score', 0) for r in all_results)
        percentage = round((total_score / max_total * 100) if max_total > 0 else 0, 1)
        
        # 收集薄弱知识点
        wrong_question_ids = [str(r.get('question_id', '')) for r in all_results if not r.get('is_correct', False)]
        weakness_tags = [r.get('knowledge_tag', '') for r in all_results if r.get('weakness') and r.get('weakness') != []]
        weakness_set = set([t for t in weakness_tags if t])
        
        # 生成改进建议
        improvement_prompt = f"""基于以下薄弱知识点：{json.dumps(list(weakness_set), ensure_ascii=False)}

请生成3-5条针对性的学习改进建议，用中文输出。"""
        
        try:
            response = llm.invoke(improvement_prompt)
            suggestions_data = parse_json_response(response.content)
            if suggestions_data:
                improvement_suggestions = suggestions_data.get('suggestions', suggestions_data.get('improvement', []))
                if isinstance(improvement_suggestions, str):
                    improvement_suggestions = [improvement_suggestions]
            else:
                improvement_suggestions = [
                    "建议回顾相关知识点，加强理解",
                    "多做相关练习题巩固知识点",
                    "注意审题，仔细阅读题目要求"
                ]
        except:
            improvement_suggestions = [
                "建议回顾相关知识点，加强理解",
                "多做相关练习题巩固知识点",
                "注意审题，仔细阅读题目要求"
            ]

        # 生成总体评价
        summary_parts = []
        if percentage >= 90:
            summary_parts.append("优秀")
        elif percentage >= 70:
            summary_parts.append("良好")
        elif percentage >= 60:
            summary_parts.append("及格")
        else:
            summary_parts.append("需要加强学习")
        summary_parts.append(f"得分 {total_score}/{max_total}")
        
        summary = "，".join(summary_parts)

        grading_result = {
            "student_id": "student_" + str(int(__import__("time").time())),
            "exam_id": "exam_" + str(int(__import__("time").time())),
            "total_score": total_score,
            "max_total": max_total,
            "percentage": percentage,
            "summary": summary,
            "details": all_results,
            "wrong_question_ids": wrong_question_ids,
            "weakness_summary": "、".join(list(weakness_set)[:5]) if weakness_set else "无明显薄弱点",
            "improvement_suggestions": improvement_suggestions
        }

        logger.info(f"Grading complete: {total_score}/{max_total}, {percentage}%")
        return {"status": "success", "grading": grading_result}
    except Exception as e:
        logger.error(f"grade_exam error: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"批改失败: {str(e)}")

exam_answers_store = {}

class SaveAnswersRequest(BaseModel):
    session_id: str
    questions: list
    answers: list

class ExplainRequest(BaseModel):
    questions: list
    answers: list

@router.post("/exam/save-answers")
async def save_answers(request: SaveAnswersRequest):
    """保存用户答案（暂存）"""
    exam_answers_store[request.session_id] = {
        "questions": request.questions,
        "answers": request.answers,
        "timestamp": __import__("datetime").datetime.now().isoformat()
    }
    answered_count = sum(1 for a in request.answers if a.get("answer", "").strip())
    return {
        "status": "success",
        "message": f"已保存 {answered_count}/{len(request.questions)} 题",
        "saved_count": answered_count,
        "total_count": len(request.questions)
    }

@router.get("/exam/get-answers/{session_id}")
async def get_answers(session_id: str):
    """获取已保存的答案"""
    if session_id not in exam_answers_store:
        return {"status": "success", "found": False, "message": "暂无保存的答案"}
    data = exam_answers_store[session_id]
    return {
        "status": "success",
        "found": True,
        "questions": data["questions"],
        "answers": data["answers"],
        "timestamp": data["timestamp"]
    }

@router.post("/exam/explain")
async def explain_answers(request: ExplainRequest):
    """获取每道题的详细AI解析"""
    try:
        from app.core.config import settings
        from langchain_openai import ChatOpenAI

        _is_siliconflow = "siliconflow" in settings.LLM_API_BASE.lower()
        _no_think_params = (
            {"enable_thinking": False}
            if _is_siliconflow
            else {"chat_template_kwargs": {"enable_thinking": False}}
        )

        llm = ChatOpenAI(
            openai_api_base=settings.LLM_API_BASE,
            openai_api_key=settings.LLM_API_KEY,
            model_name=settings.LLM_MODEL_NAME,
            temperature=0.3,
            max_tokens=4096,
            extra_body={
                **_no_think_params,
                "stop": ["<|im_end|>", "<|endoftext|>"]
            }
        )

        explain_prompt = f"""请为以下考题提供详细的AI解析。

原题：
{json.dumps(request.questions, ensure_ascii=False, indent=2)}

学生答案：
{json.dumps(request.answers, ensure_ascii=False, indent=2)}

请以JSON格式输出，包含以下字段：
- "explanations": [
    {{
        "question_id": 1,
        "question": "题目内容",
        "correct_answer": "正确答案",
        "student_answer": "学生答案",
        "analysis": "详细解析",
        "key_points": ["关键要点1", "关键要点2"],
        "related_knowledge": "相关知识点"
    }}
  ]

请确保输出是有效的JSON格式。"""

        logger.info(f"Generating explanations for {len(request.questions)} questions")
        response = llm.invoke(explain_prompt)
        logger.info(f"LLM response length: {len(response.content)}")

        try:
            content = response.content.strip()
            if content.startswith("```"):
                content = re.sub(r'^```\w*\n?', '', content)
                content = re.sub(r'\n?```$', '', content)
            json_match = re.search(r'\{[\s\S]*\}', content)
            if json_match:
                explain_result = json.loads(json_match.group())
            else:
                explain_result = {"error": "JSON解析失败", "raw": response.content[:500]}
        except json.JSONDecodeError as je:
            logger.error(f"JSON decode error: {je}")
            explain_result = {"error": f"JSON解析失败: {str(je)}", "raw": response.content[:500]}

        return {"status": "success", "explanations": explain_result}
    except Exception as e:
        logger.error(f"explain_answers error: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"生成解析失败: {str(e)}")
